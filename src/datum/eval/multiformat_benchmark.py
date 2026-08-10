"""Multi-format ingestion benchmark (task #30).

Proves the claim the DoclingParser exists to make: one file per format family,
poured through the SAME write path -> derivation -> hybrid retrieval, and the
fact each file carries is retrievable by a semantic query. Plain text/markdown
go through the dependency-free MarkdownParser (`Corpus.ingest`); every other
format goes through Docling (`Corpus.ingest_file`) — same records, same views,
same operators, so a passing benchmark is evidence that format coverage is a
write-path feeder and not a change to retrieval.

**Coverage is environment-honest, not aspirational.** The formats exercised
here are the ones whose Docling backend needs no downloaded model:
md, txt, html, csv, docx, pptx, xlsx. Three families are deliberately NOT run
here and are reported as skipped-with-reason rather than silently omitted:
  - PDF and IMAGE/scanned — need Docling's layout (and OCR) models, which
    download from HuggingFace on first use; that egress is unavailable in this
    environment (the models are not cached and the network HEAD calls hang).
    The DoclingParser handles them unchanged once the models are staged.
  - AUDIO — the Whisper model IS cached, but there is no honest way to
    synthesize speech with a KNOWN transcript to assert retrieval against
    without a TTS, so a fabricated-content audio case is worse than none.

Generating docx/pptx/xlsx needs python-docx / python-pptx / openpyxl (Docling
dependencies, so present wherever the parse extra is). A format whose generator
library is missing is reported skipped, never a benchmark crash.
"""

from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from typing import Callable

from datum.corpus import Corpus
from datum.kernel.principal import Principal

_NAMESPACE = "bench"


@dataclass(frozen=True)
class FormatCase:
    fmt: str
    filename: str
    binary: bool  # True -> Corpus.ingest_file (Docling); False -> Corpus.ingest (text)
    write: Callable[[Path], None]
    query: str
    expect_substring: str


@dataclass(frozen=True)
class FormatResult:
    fmt: str
    ingested: bool
    retrieved: bool
    detail: str


@dataclass(frozen=True)
class BenchmarkReport:
    results: tuple[FormatResult, ...]
    skipped: tuple[tuple[str, str], ...]  # (format, reason)

    @property
    def passed(self) -> bool:
        return bool(self.results) and all(r.ingested and r.retrieved for r in self.results)


# --- file generators (one known fact per format, distinct topics so each
#     query retrieves its OWN document) ---


def _w_text(text: str) -> Callable[[Path], None]:
    return lambda p: p.write_text(text, encoding="utf-8")


def _w_docx(p: Path) -> None:
    import docx

    d = docx.Document()
    d.add_heading("Access Review", level=1)
    d.add_paragraph("The quarterly audit reviews every privileged access grant and revokes stale ones.")
    d.save(p)


def _w_pptx(p: Path) -> None:
    import pptx

    pres = pptx.Presentation()
    slide = pres.slides.add_slide(pres.slide_layouts[1])
    slide.shapes.title.text = "Endpoint Security"
    slide.placeholders[1].text = "Every company laptop enforces full disk encryption while at rest."
    pres.save(p)


def _w_xlsx(p: Path) -> None:
    import openpyxl

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.append(["control", "requirement"])
    ws.append(["password rotation", "passwords must be rotated every sixty days"])
    ws.append(["mfa", "hardware security keys are mandatory for admins"])
    wb.save(p)


def _cases() -> list[FormatCase]:
    return [
        FormatCase("md", "deploy.md", False,
                   _w_text("# Deploy\n\nThe canary rollout waits thirty minutes before promoting to full production.\n"),
                   "how long does the canary wait before full rollout", "thirty minutes"),
        FormatCase("txt", "backups.txt", False,
                   _w_text("The nightly backup retains snapshots for ninety days before they are pruned.\n"),
                   "how long are nightly backup snapshots kept", "ninety days"),
        FormatCase("html", "support.html", True,
                   _w_text("<html><body><h1>Support SLA</h1><p>Premium customers receive a first response within four business hours.</p></body></html>"),
                   "how quickly do premium customers get a first response", "four business hours"),
        FormatCase("csv", "refunds.csv", True,
                   _w_text("policy,value\nrefund window,customers may request a refund within fourteen days of purchase\n"),
                   "how long is the refund window for a purchase", "fourteen days"),
        FormatCase("docx", "access_review.docx", True, _w_docx,
                   "what does the quarterly audit review", "privileged access"),
        FormatCase("pptx", "endpoint_security.pptx", True, _w_pptx,
                   "what encryption do company laptops use", "full disk encryption"),
        FormatCase("xlsx", "controls.xlsx", True, _w_xlsx,
                   "how often must passwords be rotated", "sixty days"),
    ]


_SKIPPED = (
    ("pdf", "needs Docling layout models (HuggingFace egress unavailable / not cached here)"),
    ("image/scanned", "needs Docling layout + OCR models (HuggingFace egress unavailable here)"),
    ("audio", "Whisper is cached, but no honest known-transcript audio can be synthesized without a TTS"),
)


def run_benchmark(corpus: Corpus, workdir: Path) -> BenchmarkReport:
    """Generate, ingest, and retrieve one file per working format. Returns a
    per-format report; a generator-library-missing format is skipped, not
    fatal. The caller owns `corpus` (point it at a scratch DB) and `workdir`.
    """
    principal = Principal(id="bench", namespace=_NAMESPACE)
    results: list[FormatResult] = []
    skipped = list(_SKIPPED)

    for case in _cases():
        path = workdir / case.filename
        try:
            case.write(path)
        except ImportError as exc:
            skipped.append((case.fmt, f"generator library unavailable: {exc}"))
            continue

        try:
            if case.binary:
                corpus.ingest_file(str(path), principal, source_id=case.fmt)
            else:
                corpus.ingest(case.fmt, path.read_text(encoding="utf-8"), principal=principal)
            ingested = True
        except Exception as exc:  # a real ingest failure is a benchmark FAIL, reported
            results.append(FormatResult(case.fmt, False, False, f"ingest raised {type(exc).__name__}: {exc}"))
            continue

        evidence = corpus.search(case.query, principal=principal)
        combined = "\n".join(h.content for h in evidence.hits)
        retrieved = case.expect_substring in combined
        detail = (
            f"query {case.query!r} -> status={evidence.status}, {len(evidence.hits)} hit(s); "
            + ("found expected content" if retrieved else f"MISSING {case.expect_substring!r}")
        )
        results.append(FormatResult(case.fmt, ingested, retrieved, detail))

    return BenchmarkReport(results=tuple(results), skipped=tuple(skipped))
